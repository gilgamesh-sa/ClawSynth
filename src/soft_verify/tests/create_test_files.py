#!/usr/bin/env python3
"""Generate test fixture files for soft_tools testing.

Run this script once to create all necessary test files in tests/files/.
"""
from __future__ import annotations

import csv
import json
import sys
from pathlib import Path

FILES_DIR = Path(__file__).resolve().parent / "files"


def create_text_file():
    (FILES_DIR / "sample.txt").write_text(
        "Hello, this is a test file.\n第二行中文内容。\nThird line.\n",
        encoding="utf-8",
    )


def create_json_file():
    data = {
        "name": "test",
        "items": [1, 2, 3],
        "nested": {"key": "value", "list": [{"a": 1}, {"b": 2}]},
    }
    (FILES_DIR / "sample.json").write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def create_csv_file():
    with (FILES_DIR / "sample.csv").open("w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["name", "age", "city"])
        writer.writerow(["Alice", "30", "Beijing"])
        writer.writerow(["Bob", "25", "Shanghai"])
        writer.writerow(["Carol", "28", "Guangzhou"])


def create_html_file():
    (FILES_DIR / "sample.html").write_text(
        """<!DOCTYPE html>
<html>
<head><title>Test Page</title></head>
<body>
<h1>Hello World</h1>
<p>This is a test HTML page.</p>
<table>
<tr><th>Name</th><th>Value</th></tr>
<tr><td>Alpha</td><td>100</td></tr>
</table>
</body>
</html>
""",
        encoding="utf-8",
    )


def create_excel_file():
    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["name", "score", "grade"])
        ws.append(["Alice", 95, "A"])
        ws.append(["Bob", 82, "B"])
        ws.append(["Carol", 78, "C"])
        wb.save(FILES_DIR / "sample.xlsx")
    except ImportError:
        print("SKIP: openpyxl not installed, cannot create sample.xlsx")


def create_pptx_file():
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        slide.placeholders[1].text = "This is a test presentation."
        prs.save(FILES_DIR / "sample.pptx")
    except ImportError:
        print("SKIP: python-pptx not installed, cannot create sample.pptx")


def create_image_file():
    try:
        from PIL import Image

        img = Image.new("RGB", (200, 100), color=(255, 128, 0))
        img.save(FILES_DIR / "sample.png")
    except ImportError:
        print("SKIP: Pillow not installed, cannot create sample.png")


def create_pdf_file():
    # A minimal valid PDF containing the text "hello"
    minimal_pdf = b'''%PDF-1.1
%\xa7\xf3\xa0\x8d
1 0 obj
<< /Type /Catalog /Outlines 2 0 R /Pages 3 0 R >>
endobj
2 0 obj
<< /Type /Outlines /Count 0 >>
endobj
3 0 obj
<< /Type /Pages /Kids [4 0 R] /Count 1 >>
endobj
4 0 obj
<< /Type /Page /Parent 3 0 R /MediaBox [0 0 612 792] /Contents 5 0 R /Resources << /Font << /F1 6 0 R >> >> >>
endobj
5 0 obj
<< /Length 44 >>
stream
BT
/F1 24 Tf
100 700 Td
(hello) Tj
ET
endstream
endobj
6 0 obj
<< /Type /Font /Subtype /Type1 /Name /F1 /BaseFont /Helvetica >>
endobj
xref
0 7
0000000000 65535 f 
0000000018 00000 n 
0000000077 00000 n 
0000000123 00000 n 
0000000180 00000 n 
0000000299 00000 n 
0000000394 00000 n 
trailer
<< /Size 7 /Root 1 0 R >>
startxref
482
%%EOF
'''
    (FILES_DIR / "sample.pdf").write_bytes(minimal_pdf)


def create_markdown_file():
    (FILES_DIR / "sample.md").write_text(
        "# Test Document\n\n"
        "## Section 1\n\n"
        "This is a **test** markdown file.\n\n"
        "- Item 1\n"
        "- Item 2\n",
        encoding="utf-8",
    )


def main():
    FILES_DIR.mkdir(parents=True, exist_ok=True)
    creators = [
        ("sample.txt", create_text_file),
        ("sample.json", create_json_file),
        ("sample.csv", create_csv_file),
        ("sample.html", create_html_file),
        ("sample.xlsx", create_excel_file),
        ("sample.pptx", create_pptx_file),
        ("sample.png", create_image_file),
        ("sample.pdf", create_pdf_file),
        ("sample.md", create_markdown_file),
    ]
    for name, fn in creators:
        fn()
        path = FILES_DIR / name
        status = "OK" if path.exists() else "MISSING"
        print(f"  {status}: {name}")

    print(f"\nFiles created in: {FILES_DIR}")


if __name__ == "__main__":
    main()
